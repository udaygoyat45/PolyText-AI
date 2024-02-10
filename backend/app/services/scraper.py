from app import db
import celery
import pptx
from pptx import Presentation
import fitz
from PIL import Image
from celery import shared_task
import io

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    text = []

    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)

    text = " ".join(text)
    return text


def extract_text_from_pdf(pdf_path):
    text = ''
    # Open the PDF file
    with fitz.open(pdf_path) as pdf_document:
        # Iterate through each page
        for page_number in range(pdf_document.page_count):
            # Get the page
            page = pdf_document.load_page(page_number)

            # Extract text from the page
            text += page.get_text()

    return text


def extract_text_from_text(text_path):
    text = ''
    with open(text_path) as fin:
        text = ''.join(fin.readlines())
    
    return text
        


@shared_task()
def scrape(file_id):
    media = db.media_sources.find_one({'id': file_id})
    if media['scraped']:
        return

    match media['type']:
        case 'application/pdf':
            pdf_text = extract_text_from_pdf(media['location'])
            db.media_sources.update_one({'id': file_id}, {'$set': {'scraped': True, 'scraped_text': pdf_text}})

        case 'application-vnd.openxmlformats-officedocument.presentationml.presentation':
            pptx_text = extract_text_from_pptx(media['location'])
            db.media_sources.update_one({'id': file_id}, {'$set': {'scraped': True, 'scraped_text': pptx_text}})
        
        case 'plain/text':
            text_text = extract_text_from_text(media['location'])
            db.media_sources.update_one({'id': file_id}, {'$set': {'scraped': True, 'scraped_text': text_text}})
